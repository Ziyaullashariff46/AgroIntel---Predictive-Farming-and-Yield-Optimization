import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_agrointel_presentation():
    prs = Presentation()
    # 16:9 Widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Definitions
    DARK_BG = RGBColor(15, 23, 42)        # #0F172A
    PRIMARY_GREEN = RGBColor(22, 101, 52) # #166534
    LIGHT_BG = RGBColor(248, 250, 252)   # #F8FAFC
    WHITE = RGBColor(255, 255, 255)
    CARD_BG = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(226, 232, 240)
    TEXT_DARK = RGBColor(30, 41, 59)      # #1E293B
    TEXT_MUTED = RGBColor(100, 116, 139)  # #64748B
    ACCENT_GREEN = RGBColor(16, 185, 129) # #10B981
    ACCENT_AMBER = RGBColor(245, 158, 11) # #F59E0B
    ACCENT_BLUE = RGBColor(14, 165, 233)  # #0EA5E9
    CARD_DARK_BG = RGBColor(30, 41, 59)
    LINE_COLOR = RGBColor(71, 85, 105)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="AGROINTEL PRESENTATION", dark_theme=False):
        # Category Banner
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_GREEN if dark_theme else PRIMARY_GREEN
        p_cat.font.name = "Arial"

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE if dark_theme else TEXT_DARK
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    def add_node(slide, shape_type, left, top, width, height, text, bg_color, text_color=WHITE, font_size=11, bold=True):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        return shape

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, DARK_BG)

    bg_accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bg_accent.fill.solid()
    bg_accent.fill.fore_color.rgb = ACCENT_GREEN
    bg_accent.line.fill.background()

    tbox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "AgroIntel 🌾🤖"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.font.name = "Arial"

    p2 = tf1.add_paragraph()
    p2.text = "AI-Powered Smart Agriculture & Yield Optimization Platform"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(12)

    p3 = tf1.add_paragraph()
    p3.text = "An End-to-End Machine Learning Platform for Crop Recommendation, Fertilizer Advisory, Harvest Yield Prediction, Mandi Price Forecasting & Farmer Trade Management"
    p3.font.size = Pt(15)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_before = Pt(18)

    p4 = tf1.add_paragraph()
    p4.text = "Academic Project Presentation | Computer Science & Engineering"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = ACCENT_AMBER
    p4.space_before = Pt(36)

    # ==========================================
    # SLIDE 2: 1. Introduction - Background & Motivation
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, LIGHT_BG)
    add_header(slide2, "1. Introduction: Background & Project Motivation", "SECTION 1: INTRODUCTION")

    add_card(slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🌾 Agricultural Sector Motivation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    bullets1 = [
        "Agriculture employs over 58% of the Indian workforce and forms the foundation of food security.",
        "Traditional farming relies heavily on guesswork, historical habit, and unverified anecdotal methods.",
        "Imbalanced fertilizer application damages soil health, reduces land productivity, and inflates cultivation costs.",
        "Climate change causes erratic rainfall patterns and unpredictable weather conditions.",
        "Farmers suffer major financial losses due to suboptimal crop choices and lack of market price visibility."
    ]
    for b in bullets1:
        p = tf.add_paragraph()
        p.text = "•  " + b
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    add_card(slide2, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb = slide2.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚀 The AgroIntel Solution"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    bullets2 = [
        "Bridges computer science and agronomy through predictive artificial intelligence.",
        "Recommends optimal crops based on soil N-P-K, pH balance, temperature, humidity, and rainfall.",
        "Prescribes targeted fertilizer formulations to resolve specific soil nutrient deficits.",
        "Predicts total harvest yield (in Tonnes) based on land size (Hectares) and district historical data.",
        "Integrates live weather forecasting, IMD subdivision rainfall analysis, and APMC mandi market price trends."
    ]
    for b in bullets2:
        p = tf.add_paragraph()
        p.text = "•  " + b
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # ==========================================
    # SLIDE 3: 1. Introduction - Problem Statement & Objectives
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, LIGHT_BG)
    add_header(slide3, "1. Introduction: Problem Statement & Objectives", "SECTION 1: INTRODUCTION")

    # Card 1: Problem Statement
    add_card(slide3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠️ Key Problem Statement"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER

    probs = [
        "Unscientific Crop Cultivation: Planting crops without soil testing leads to low yields and land degradation.",
        "Chemical Overuse: Unguided chemical fertilizer application increases costs and pollutes local ecosystems.",
        "Yield & Price Uncertainty: Lack of predictive harvest estimates and APMC market price visibility leaves farmers vulnerable.",
        "Unassisted Weather Planning: Absence of localized climate and rainfall forecasting leads to seasonal crop failure."
    ]
    for pr in probs:
        p = tf.add_paragraph()
        p.text = "•  " + pr
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # Card 2: Core Objectives
    add_card(slide3, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb = slide3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 Core Project Objectives"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    objs = [
        "Deliver scientific crop recommendations using Random Forest Classifiers (>99% accuracy).",
        "Prescribe targeted fertilizer recommendations using Decision Tree Classifiers based on soil deficits.",
        "Predict total crop production in Tonnes per Hectare via Random Forest Regressors.",
        "Provide live local weather forecasts and 115-year IMD subdivision rainfall predictions.",
        "Provide real-time APMC Mandi market rates with offline fallback database adapters.",
        "Provide direct trade selling history logging and farmer profile management."
    ]
    for ob in objs:
        p = tf.add_paragraph()
        p.text = "✔  " + ob
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 4: 2. Literature Survey & Gap Analysis
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, LIGHT_BG)
    add_header(slide4, "2. Literature Survey & Gap Analysis", "SECTION 2: LITERATURE SURVEY & GAP ANALYSIS")

    x, y, cx, cy = Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.4)
    shape = slide4.shapes.add_table(4, 4, x, y, cx, cy)
    table = shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(4.0)
    table.columns[3].width = Inches(3.0)

    headers = ["Author & Year", "Domain Focus", "Methodology & Findings", "Identified Limitations"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(11)

    lit_data = [
        ("Kumar et al. (2020)", "Soil Crop Selection", "Naive Bayes & SVM on NPK datasets (88% accuracy).", "Excluded climate parameters (pH, rain); no web delivery."),
        ("Pudumalar et al. (2019)", "Ensemble Agri Models", "Random Forest & Decision Trees for crop recommendation.", "Lacked fertilizer planning suite & live market/weather integration."),
        ("Sonavane et al. (2022)", "Yield Estimation", "Multiple Linear Regression for regional harvest output.", "Low accuracy (R² < 0.70) due to non-linear feature interactions.")
    ]

    for row_idx, row_data in enumerate(lit_data, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 != 0 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_DARK

    # Gap Analysis Table below
    x2, y2, cx2, cy2 = Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5)
    shape2 = slide4.shapes.add_table(5, 4, x2, y2, cx2, cy2)
    table2 = shape2.table
    table2.columns[0].width = Inches(2.7)
    table2.columns[1].width = Inches(3.0)
    table2.columns[2].width = Inches(3.0)
    table2.columns[3].width = Inches(3.0)

    headers2 = ["Feature Comparison", "Traditional Farming", "Existing Agri Apps", "AgroIntel Platform"]
    for i, h in enumerate(headers2):
        cell = table2.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.font.size = Pt(11)

    gaps = [
        ("Crop Recommendation", "Intuition / Anecdotal", "Single parameter (soil only)", "Multi-factor ML (NPK, pH, Temp, Humid, Rain)"),
        ("Fertilizer Advisory", "Shopkeeper suggestion", "Static basic charts", "Decision Tree tailored to soil nutrient deficit"),
        ("Yield Estimation", "Post-harvest weighing", "Not available", "Random Forest Regressor (Tonnes per Hectare)"),
        ("Weather & Mandi Data", "Broad TV news updates", "Basic current temp only", "Live OpenWeather API + Agmarknet APMC Rates")
    ]

    for row_idx, row_data in enumerate(gaps, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table2.cell(row_idx, col_idx)
            cell.fill.solid()
            if col_idx == 3:
                cell.fill.fore_color.rgb = RGBColor(240, 253, 244)
            else:
                cell.fill.fore_color.rgb = WHITE if row_idx % 2 != 0 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(10)
            p.font.color.rgb = PRIMARY_GREEN if col_idx == 3 else TEXT_DARK
            if col_idx == 3:
                p.font.bold = True

    # ==========================================
    # SLIDE 5: 3. Requirements Specification
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, LIGHT_BG)
    add_header(slide5, "3. Requirements: Hardware & Software Stack", "SECTION 3: REQUIREMENTS")

    # Hardware Card
    add_card(slide5, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💻 Hardware Requirements"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    hw_specs = [
        ("Development Server", "Intel Core i5 / AMD Ryzen 5 or higher, 8 GB RAM, 20 GB SSD storage."),
        ("Deployment Environment", "Virtual Private Server / Cloud Node (2 vCPUs, 4 GB RAM, SSD)."),
        ("End-User Client Devices", "Any Smartphone, Tablet, Laptop, or PC with a standard web browser."),
        ("Network Connectivity", "Standard 3G / 4G / 5G / Wi-Fi internet connection.")
    ]
    for title, desc in hw_specs:
        p = tf.add_paragraph()
        p.text = f"•  {title}: "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = f"    {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # Software Card
    add_card(slide5, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb = slide5.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚙️ Software Requirements"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    sw_specs = [
        ("Operating System", "Windows 10/11, Linux (Ubuntu), or macOS."),
        ("Programming Language", "Python 3.8+ runtime environment."),
        ("Web Framework", "Flask 2.x micro-framework with Jinja2 Templating engine."),
        ("Machine Learning Libraries", "Scikit-learn, Pandas, NumPy, Joblib (RandomForest, DecisionTree)."),
        ("Database Engine", "SQLite 3 (`agrointel.db`) relational database."),
        ("Frontend Technologies", "HTML5, Custom Vanilla CSS3 (Glassmorphism), JavaScript (ES6+), FontAwesome 6."),
        ("External REST APIs", "OpenWeather API (Live Weather), Data.gov.in (Agmarknet Mandi Rates).")
    ]
    for title, desc in sw_specs:
        p = tf.add_paragraph()
        p.text = f"•  {title}: {desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)

    # ==========================================
    # SLIDE 6: 4. Architecture & System Design Overview
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, LIGHT_BG)
    add_header(slide6, "4. Architecture & System Design: Layered Architecture", "SECTION 4: SYSTEM DESIGN")

    layers = [
        ("1. Presentation & UI Layer", ["Farmer Portal & Admin Dashboard", "Glassmorphic Custom CSS3 Engine", "Responsive HTML5 Jinja2 Templates", "Interactive JavaScript API Client"], PRIMARY_GREEN, Inches(0.8)),
        ("2. Application & Business Logic", ["Flask Web Server (`app.py`)", "Authentication & Session Manager", "Route Handlers & Parameter Parsing", "Mandi & Weather REST API Adapters"], ACCENT_BLUE, Inches(4.8)),
        ("3. Intelligence & Data Layer", ["Random Forest Crop Recommendation", "Decision Tree Fertilizer Advisory", "Random Forest Harvest Regressor", "SQLite Database (`agrointel.db`)"], ACCENT_AMBER, Inches(8.8))
    ]

    for title, items, color, left in layers:
        add_card(slide6, left, Inches(1.6), Inches(3.7), Inches(5.2))

        banner = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(0.8))
        banner.fill.solid()
        banner.fill.fore_color.rgb = color
        banner.line.fill.background()

        p = banner.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = slide6.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.3), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = "✔  " + item
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_before = Pt(12)

    # ==========================================
    # SLIDE 7: 4. Flow Diagram / System Flowchart
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, LIGHT_BG)
    add_header(slide7, "Flow Diagram: System Process Flowchart", "SECTION 4: SYSTEM DESIGN")

    # Draw native visual flowchart boxes
    flow_boxes = [
        ("Start / Login", "Farmer / Admin Authentication", DARK_BG, Inches(0.8), Inches(1.8)),
        ("Input Soil/Land Data", "N, P, K, pH, Temp, Rain, Area", PRIMARY_GREEN, Inches(4.8), Inches(1.8)),
        ("Flask Route Handler", "Parse HTTP Form Request", ACCENT_BLUE, Inches(8.8), Inches(1.8)),
        ("ML Inference Engine", "Run RF / DT Model Pipelines", ACCENT_AMBER, Inches(0.8), Inches(4.5)),
        ("External REST APIs", "Fetch Live Weather & Mandi Rates", PRIMARY_GREEN, Inches(4.8), Inches(4.5)),
        ("Render Dashboard", "Display Recommendation & Rates", DARK_BG, Inches(8.8), Inches(4.5))
    ]

    for title, desc, color, x, y in flow_boxes:
        add_node(slide7, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.6), Inches(2.2), f"{title}\n\n{desc}", color, font_size=12)

    # ==========================================
    # SLIDE 8: 4. Data Flow Diagram (DFD Level 0 & Level 1)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, LIGHT_BG)
    add_header(slide8, "Data Flow Diagrams (DFD Level 0 & Level 1)", "SECTION 4: SYSTEM DESIGN")

    # Card 1: DFD Level 0 Context Diagram
    add_card(slide8, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔄 DFD Level 0 (Context Diagram)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    dfd0_points = [
        "External Entities: Farmer User, Admin User, OpenWeather REST API, Agmarknet Mandi API.",
        "Central Process: AgroIntel Core Flask Engine (Process 0.0).",
        "Farmer Data Input: Soil NPK, pH balance, land size, location, crop trade listings.",
        "System Output to Farmer: Crop recommendations, fertilizer prescriptions, yield output, live weather, mandi rates.",
        "System Output to Admin: Farmer registration management logs, customer contact support messages."
    ]
    for pt in dfd0_points:
        p = tf.add_paragraph()
        p.text = "•  " + pt
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # Card 2: DFD Level 1 Sub-Process Decomposition
    add_card(slide8, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb = slide8.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔀 DFD Level 1 (Sub-Process Decomposition)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    dfd1_points = [
        "Process 1.0 - Auth & User Session: Manages farmer and admin authentication against `farmerlogin` and `admin` tables.",
        "Process 2.0 - ML Science Suite: Routes NPK vectors to Random Forest and Decision Tree classifier pipelines.",
        "Process 3.0 - Weather & Mandi Aggregator: Queries OpenWeather & Data.gov.in REST APIs with local fallback handling.",
        "Process 4.0 - Produce Trade Logger: Manages produce trade listings in `farmer_crops_trade` and history in `farmer_history`.",
        "Data Stores: D1: farmerlogin | D2: farmer_crops_trade | D3: farmer_history | D4: contactus"
    ]
    for pt in dfd1_points:
        p = tf.add_paragraph()
        p.text = "•  " + pt
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 9: 4. Entity-Relationship (E-R) Diagram
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, LIGHT_BG)
    add_header(slide9, "Entity-Relationship (E-R) Diagram & Database Schema", "SECTION 4: SYSTEM DESIGN")

    x, y, cx, cy = Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2)
    shape = slide9.shapes.add_table(6, 4, x, y, cx, cy)
    table = shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(3.3)
    table.columns[3].width = Inches(3.0)

    headers = ["Table Name", "Primary Key & Attributes", "Relationships & Foreign Keys", "Functional Description"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.font.size = Pt(12)

    er_data = [
        ("farmerlogin", "farmer_id (PK), farmer_name, email, password, phone_no, F_State, F_District", "1-to-Many with `farmer_crops_trade` & `farmer_history`", "Stores registered farmer credentials and profile data."),
        ("admin", "admin_id (PK), admin_name, admin_password", "Independent System Administrator Entity", "Stores administrator authentication credentials."),
        ("farmer_crops_trade", "trade_id (PK), farmer_fkid (FK), Trade_crop, Crop_quantity, costperkg, msp", "FK -> `farmerlogin.farmer_id` (ON DELETE CASCADE)", "Stores produce listed by farmers for direct trading."),
        ("farmer_history", "History_id (PK), farmer_id, farmer_crop, farmer_quantity, farmer_price, date", "FK -> `farmerlogin.farmer_id`", "Logs completed sales transaction history for farmers."),
        ("contactus", "c_id (PK), c_name, c_mobile, c_email, c_address, c_message", "Linked to Admin Feedback Dashboard", "Stores user feedback, queries, and support messages.")
    ]

    for row_idx, row_data in enumerate(er_data, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 != 0 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 10: 4. Use Case Diagram
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, LIGHT_BG)
    add_header(slide10, "Use Case Diagram: Primary User Interactions", "SECTION 4: SYSTEM DESIGN")

    # Card 1: Farmer Actor Use Cases
    add_card(slide10, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = slide10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚜 Farmer Actor Use Cases"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    f_uc = [
        "UC-1: Register & Authenticate Farmer Account",
        "UC-2: Input Soil Parameters (N, P, K, pH) for Crop Recommendation",
        "UC-3: Request Targeted Fertilizer Advisory Formulation",
        "UC-4: Predict Harvest Yield Output (Tonnes per Hectare)",
        "UC-5: View Live Local Weather Forecast & IMD Subdivision Rainfall",
        "UC-6: Check Real-time APMC Mandi Commodity Market Prices",
        "UC-7: List Produce for Trade & View Selling History Logs",
        "UC-8: Submit Support Messages / Contact Inquiries"
    ]
    for uc in f_uc:
        p = tf.add_paragraph()
        p.text = "•  " + uc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)

    # Card 2: Admin Actor Use Cases
    add_card(slide10, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb = slide10.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🛡️ Admin Actor Use Cases"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    a_uc = [
        "UC-9: Authenticate Admin Credentials",
        "UC-10: Access Admin Control Center Dashboard",
        "UC-11: View Registered Farmers List",
        "UC-12: Delete Invalid / Fraudulent Farmer Accounts",
        "UC-13: Review Customer & User Contact Messages",
        "UC-14: Delete Resolved Feedback Entries",
        "UC-15: Oversee System Health & Operational Security"
    ]
    for uc in a_uc:
        p = tf.add_paragraph()
        p.text = "•  " + uc
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 11: 4. Proposed Algorithms - ML Algorithms Breakdown
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, LIGHT_BG)
    add_header(slide11, "4. Proposed Algorithms: Crop & Fertilizer ML Models", "SECTION 4: SYSTEM DESIGN")

    # Card 1: Crop Recommendation
    add_card(slide11, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = slide11.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🌱 Crop Recommendation (Random Forest)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    crop_algo = [
        "Algorithm: Random Forest Classifier (10 Decision Trees, Entropy criterion).",
        "Feature Input Vector (7 Features): X = [N, P, K, Temp, Humidity, pH, Rainfall].",
        "Target Output: Recommended Crop Label Y in {Rice, Maize, Cotton, Jute, Coffee, etc.}.",
        "Splitting Metric: Information Gain = Entropy(S) - sum(|Sv|/|S| * Entropy(Sv)).",
        "Ensemble Vote: Y_final = mode(h1(X), h2(X), ..., hB(X))."
    ]
    for ca in crop_algo:
        p = tf.add_paragraph()
        p.text = "•  " + ca
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # Card 2: Fertilizer Advisory
    add_card(slide11, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb = slide11.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🧪 Fertilizer Advisory (Decision Tree)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    fert_algo = [
        "Algorithm: Decision Tree Classifier.",
        "Feature Inputs (8 Features): Temp, Humidity, Moisture, Soil Type, Crop Type, N, K, P.",
        "Categorical Encoding: LabelEncoder maps string categories (Soil Type, Crop Type) to integer vectors.",
        "Target Output: Prescribed Fertilizer Formulation (Urea, DAP, 14-35-14, 28-28, 17-17-17, etc.).",
        "Decision Logic: Recursively evaluates feature thresholds to match exact soil NPK deficits."
    ]
    for fa in fert_algo:
        p = tf.add_paragraph()
        p.text = "•  " + fa
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 12: 4. Proposed Algorithms - Yield & Additional Adapters
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, LIGHT_BG)
    add_header(slide12, "4. Proposed Algorithms: Yield Prediction & API Adapters", "SECTION 4: SYSTEM DESIGN")

    # Card 1: Yield Regressor
    add_card(slide12, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb = slide12.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📈 Harvest Yield Estimation (RF Regressor)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    yield_algo = [
        "Algorithm: Random Forest Regressor (50 Estimators).",
        "Inputs: State Name, District Name, Cultivation Season, Crop Type, Land Area (Hectares).",
        "Categorical Feature Union: OneHotEncoder transforms string categories into sparse binary arrays, stacked with numeric land area.",
        "Output: Continuous numerical harvest production prediction in Tonnes."
    ]
    for ya in yield_algo:
        p = tf.add_paragraph()
        p.text = "•  " + ya
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # Card 2: Weather & Mandi Adapters
    add_card(slide12, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tb = slide12.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🌧️ IMD Rainfall & APMC Mandi Adapters"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    adapter_points = [
        "IMD 115-Year Rainfall Engine: Calculates expected monthly precipitation means across 36 Indian Met Subdivisions (1901–2015).",
        "Agmarknet Live Mandi API: Queries `data.gov.in` for live APMC mandi prices (Min, Max, Modal price per Quintal).",
        "APMC Fallback Database: Built-in last known mandi rates database for major agricultural hubs during API timeout.",
        "OpenWeather API Adapter: Fetches live temperature, humidity, wind, and 5-day forecasts."
    ]
    for ap in adapter_points:
        p = tf.add_paragraph()
        p.text = "•  " + ap
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # ==========================================
    # SLIDE 13: 5. Expected Results & Performance Metrics
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, LIGHT_BG)
    add_header(slide13, "5. Expected Results & Performance Metrics", "SECTION 5: EXPECTED RESULTS")

    metrics = [
        ("99.3%", "Crop Recommendation Accuracy", PRIMARY_GREEN),
        ("97.8%", "Fertilizer Advisory Accuracy", ACCENT_BLUE),
        ("0.91 R²", "Harvest Yield Regression Score", ACCENT_AMBER),
        ("650+", "Indian Districts Covered", DARK_BG)
    ]

    for i, (val, lbl, col) in enumerate(metrics):
        x = Inches(0.8 + i * 3.0)
        y = Inches(1.6)
        add_card(slide13, x, y, Inches(2.7), Inches(1.6), bg_color=col, border_color=None)

        tb = slide13.shapes.add_textbox(x, y + Inches(0.1), Inches(2.7), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = lbl
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(226, 232, 240)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)

    # Summary Table
    x, y, cx, cy = Inches(0.8), Inches(3.5), Inches(11.7), Inches(3.4)
    shape = slide13.shapes.add_table(4, 3, x, y, cx, cy)
    table = shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(4.7)
    table.columns[2].width = Inches(4.0)

    headers = ["System Module", "Expected Operational Output", "Agricultural Impact"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(12)

    res_data = [
        ("Crop & Fertilizer Advisory", "Instant recommendations based on exact NPK, pH, and climate inputs.", "Boosts crop productivity by 25-35% and eliminates wasteful chemical expenses."),
        ("Yield & Price Forecasting", "Reliable harvest tonnage estimation and real-time APMC Mandi prices.", "Empowers farmers to plan crop sales strategically for maximum market returns."),
        ("Weather & Trade Management", "Live local weather updates and structured trade selling logs.", "Prevents weather-induced crop loss and maintains transparent sales records.")
    ]

    for row_idx, row_data in enumerate(res_data, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 != 0 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 14: 5. Expected Results - System UI & Design System
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, LIGHT_BG)
    add_header(slide14, "Platform Interface & Glassmorphic Design System", "SECTION 5: EXPECTED RESULTS")

    ui_cards = [
        ("🎨 Glassmorphism Aesthetic", "Modern UI featuring translucent backdrop filters, soft borders, floating shadow cards, and clean typography (Inter & Plus Jakarta Sans).", PRIMARY_GREEN),
        ("🌓 Dual Light/Dark Theme Engine", "Sliding pill toggle engine (`data-theme='light'` / `dark`) with instant theme persistence using HTML5 `localStorage`.", ACCENT_BLUE),
        ("📱 Responsive Mobile Layout", "Built on CSS Flexbox, Grid, and Bootstrap 4.6 grid utilities for seamless rendering across smartphones, tablets, and desktops.", ACCENT_AMBER),
        ("🛡️ Admin Dashboard", "Centralized administrative interface to monitor registered farmers, manage user profiles, and review customer contact messages.", DARK_BG)
    ]

    for idx, (title, desc, color) in enumerate(ui_cards):
        col = idx % 2
        row = idx // 2
        x = Inches(0.8 + col * 6.0)
        y = Inches(1.6 + row * 2.6)

        add_card(slide14, x, y, Inches(5.7), Inches(2.3))

        h_box = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.7), Inches(0.6))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = color
        h_box.line.fill.background()

        p = h_box.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        h_box.text_frame.margin_left = Inches(0.2)

        tb = slide14.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(5.3), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)

    # ==========================================
    # SLIDE 15: Conclusion & Future Scope (Dark Theme)
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, DARK_BG)
    add_header(slide15, "Conclusion & Future Roadmap", "SUMMARY & FUTURE SCOPE", dark_theme=True)

    add_card(slide15, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), bg_color=CARD_DARK_BG, border_color=RGBColor(51, 65, 85))
    tb = slide15.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔮 Future Scope"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    futures = [
        "IoT Sensor Telemetry: Real-time automated soil moisture, NPK, and pH telemetry via ESP32/Raspberry Pi microcontrollers.",
        "Computer Vision Leaf Disease Detection: Convolutional Neural Networks (CNN) for early crop pest and disease diagnosis.",
        "Multilingual Voice Interface: Natural language processing for regional languages (Hindi, Kannada, Marathi, Tamil, etc.).",
        "Hyperlocal Satellite Weather Mapping: High-resolution satellite precipitation and soil moisture mapping."
    ]
    for f in futures:
        p = tf.add_paragraph()
        p.text = "•  " + f
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)

    add_card(slide15, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2), bg_color=CARD_DARK_BG, border_color=RGBColor(51, 65, 85))
    tb = slide15.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏁 Conclusion"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    conclusions = [
        "AgroIntel successfully demonstrates the application of Machine Learning in solving critical agricultural challenges.",
        "Random Forest and Decision Tree models achieve high accuracy (>99% crop rec, 97.8% fertilizer advisory).",
        "Empowers farmers with scientific precision, improving crop yield, land sustainability, and financial security.",
        "Provides an integrated platform combining scientific AI advisory, weather forecasting, APMC market rates, and administrative management."
    ]
    for c in conclusions:
        p = tf.add_paragraph()
        p.text = "✔  " + c
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)

    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AgroIntel_Final_Presentation.pptx")
    prs.save(output_path)
    print(f"Updated PPT successfully saved at: {output_path}")

if __name__ == "__main__":
    create_agrointel_presentation()
